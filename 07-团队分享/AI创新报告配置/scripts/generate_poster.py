#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI创新报告精美海报生成器
使用Ocean Gradient配色方案
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_ai_innovation_poster():
    """创建AI创新报告海报"""
    
    # 创建演示文稿（16:9宽屏）
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 添加空白幻灯片
    blank_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_layout)
    
    # Ocean Gradient配色方案
    COLOR_DEEP_BLUE = RGBColor(6, 90, 130)      # #065A82
    COLOR_CYAN = RGBColor(28, 114, 147)         # #1C7293
    COLOR_PURPLE = RGBColor(33, 41, 92)         # #21295C
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_LIGHT_GRAY = RGBColor(240, 240, 240)
    COLOR_GOLD = RGBColor(255, 193, 7)          # 高亮色
    
    # ========== 背景渐变 ==========
    background = slide.shapes.add_shape(
        1,  # 矩形
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_DEEP_BLUE
    background.line.fill.background()
    
    # ========== 标题区域 ==========
    # 主标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5),
        Inches(14), Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_frame.text = "AI创新信息采集报告"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # 日期和副标题
    from datetime import datetime
    today_display = datetime.now().strftime("%Y年%m月%d日")
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.8),
        Inches(14), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"{today_display} | 极客四维过滤器情报分析"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = COLOR_GOLD
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # ========== Top 3 必读区域 ==========
    top3_y_start = 2.6
    
    # Top 3标题
    top3_title = slide.shapes.add_textbox(
        Inches(1), Inches(top3_y_start),
        Inches(14), Inches(0.5)
    )
    top3_title_frame = top3_title.text_frame
    top3_title_frame.text = "🏆 今日 Top 3 必读"
    top3_para = top3_title_frame.paragraphs[0]
    top3_para.font.size = Pt(36)
    top3_para.font.bold = True
    top3_para.font.color.rgb = COLOR_GOLD
    top3_para.alignment = PP_ALIGN.LEFT
    
    # Top 3内容
    top3_items = [
        {
            "num": "1",
            "title": "Anthropic完成$300B史诗级融资，估值$380B",
            "desc": "第二家2万亿级AI独角兽 • 年化收入$140B • 准备下半年IPO",
            "icon": "💰"
        },
        {
            "num": "2",
            "title": "Google DeepMind发布Aletheia数学研究Agent",
            "desc": "自主解决开放性数学难题 • Codeforces Elo 3455 • 元认知能力突破",
            "icon": "🧮"
        },
        {
            "num": "3",
            "title": "SmolLM3：3B参数碾压4B模型",
            "desc": "支持6语言+128K上下文 • 开源小模型新标杆 • 端侧部署新选择",
            "icon": "🚀"
        }
    ]
    
    item_y = top3_y_start + 0.7
    for item in top3_items:
        # 编号圆圈
        circle = slide.shapes.add_shape(
            1,  # 矩形（可以设置圆角）
            Inches(1.2), Inches(item_y),
            Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_GOLD
        circle.line.fill.background()
        
        # 编号文字
        num_frame = circle.text_frame
        num_frame.text = item["num"]
        num_para = num_frame.paragraphs[0]
        num_para.font.size = Pt(24)
        num_para.font.bold = True
        num_para.font.color.rgb = COLOR_DEEP_BLUE
        num_para.alignment = PP_ALIGN.CENTER
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(2), Inches(item_y),
            Inches(12), Inches(0.35)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"{item['icon']} {item['title']}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(22)
        title_para.font.bold = True
        title_para.font.color.rgb = COLOR_WHITE
        
        # 描述
        desc_box = slide.shapes.add_textbox(
            Inches(2), Inches(item_y + 0.35),
            Inches(12), Inches(0.3)
        )
        desc_frame = desc_box.text_frame
        desc_frame.text = item["desc"]
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(16)
        desc_para.font.color.rgb = COLOR_LIGHT_GRAY
        
        item_y += 0.9
    
    # ========== 关键数据区域 ==========
    data_y_start = 6.2
    
    # 数据标题
    data_title = slide.shapes.add_textbox(
        Inches(1), Inches(data_y_start),
        Inches(14), Inches(0.4)
    )
    data_title_frame = data_title.text_frame
    data_title_frame.text = "📊 关键数据"
    data_para = data_title_frame.paragraphs[0]
    data_para.font.size = Pt(28)
    data_para.font.bold = True
    data_para.font.color.rgb = COLOR_GOLD
    
    # 数据卡片
    data_items = [
        {"label": "融资金额", "value": "$300", "unit": "亿"},
        {"label": "新模型发布", "value": "5", "unit": "个"},
        {"label": "GitHub热门", "value": "10", "unit": "项目"},
        {"label": "Elo评分", "value": "3455", "unit": ""}
    ]
    
    card_width = 3.2
    card_height = 1.2
    card_spacing = 0.3
    card_x_start = 1
    card_y = data_y_start + 0.6
    
    for i, item in enumerate(data_items):
        card_x = card_x_start + i * (card_width + card_spacing)
        
        # 卡片背景
        card = slide.shapes.add_shape(
            1,
            Inches(card_x), Inches(card_y),
            Inches(card_width), Inches(card_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CYAN
        card.line.color.rgb = COLOR_GOLD
        card.line.width = Pt(2)
        
        # 数值
        value_box = slide.shapes.add_textbox(
            Inches(card_x), Inches(card_y + 0.2),
            Inches(card_width), Inches(0.5)
        )
        value_frame = value_box.text_frame
        value_frame.text = item["value"]
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(40)
        value_para.font.bold = True
        value_para.font.color.rgb = COLOR_WHITE
        value_para.alignment = PP_ALIGN.CENTER
        
        # 单位
        unit_box = slide.shapes.add_textbox(
            Inches(card_x), Inches(card_y + 0.7),
            Inches(card_width), Inches(0.25)
        )
        unit_frame = unit_box.text_frame
        unit_frame.text = item["unit"]
        unit_para = unit_frame.paragraphs[0]
        unit_para.font.size = Pt(18)
        unit_para.font.color.rgb = COLOR_LIGHT_GRAY
        unit_para.alignment = PP_ALIGN.CENTER
        
        # 标签
        label_box = slide.shapes.add_textbox(
            Inches(card_x), Inches(card_y + 0.9),
            Inches(card_width), Inches(0.25)
        )
        label_frame = label_box.text_frame
        label_frame.text = item["label"]
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(16)
        label_para.font.color.rgb = COLOR_GOLD
        label_para.alignment = PP_ALIGN.CENTER
    
    # ========== 页脚 ==========
    today_footer = datetime.now().strftime("%Y%m%d")
    footer = slide.shapes.add_textbox(
        Inches(1), Inches(8.3),
        Inches(14), Inches(0.4)
    )
    footer_frame = footer.text_frame
    footer_frame.text = f"极客四维过滤器 | 技术前沿度 • 落地可行性 • 工具效率比 • 行业渗透率 | 报告编号: AIR-{today_footer}"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(14)
    footer_para.font.color.rgb = COLOR_LIGHT_GRAY
    footer_para.alignment = PP_ALIGN.CENTER
    
    # 保存文件
    from datetime import datetime
    today = datetime.now().strftime("%Y-%m-%d")
    output_path = f"/data/workspace/ai_innovation_reports/AI创新报告_{today}_海报.pptx"
    prs.save(output_path)
    print(f"✅ 海报生成成功！")
    print(f"📁 文件路径: {output_path}")
    
    return output_path

if __name__ == "__main__":
    try:
        create_ai_innovation_poster()
    except Exception as e:
        print(f"❌ 生成失败: {str(e)}")
        import traceback
        traceback.print_exc()
